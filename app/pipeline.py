import base64
import concurrent.futures
import binascii
import hashlib
import json
import os
import re
import traceback
import uuid
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from threading import Lock
from typing import Any, Callable, Iterable, Optional

import requests
from pptx import Presentation
from pptx.util import Inches

from app.model_api import chat_completion_text
from app.schemas import GenerateResponse, OutlineResult, SlideOutline, SlideResult
from app.settings import Settings

ProgressCallback = Optional[Callable[[dict[str, Any]], None]]


@dataclass
class RuntimeConfig:
    text_provider: str
    text_base_url: str
    image_provider: str
    image_base_url: str
    text_api_key: str
    image_api_key: str
    text_model: str
    image_model: str


PROMPT_SUFFIX = (
    "Render as a widescreen 16:9 presentation slide (single page), "
    "zoom out slightly, center the subject, add 8-12% safe margin/padding on all sides, "
    "ensure nothing is cropped or out of frame."
)

STYLE_PROMPT_GUARD = (
    "只描述视觉风格与版式语言，禁止出现任何PPT生成参数、画幅比例、分辨率、像素、随机种子、"
    "采样步数、CFG、负面提示词、输出格式或模型参数。"
)
STYLE_PROMPT_PROHIBITION_RE = re.compile(r"(禁止|不得|不要|避免|不能|勿|严禁|不包含|不出现|不可|去除)")
STYLE_PROMPT_PARAMETER_RE = re.compile(
    r"(?i)(?:"
    r"\b(?:aspect\s*ratio|resolution|dpi|pixels?|seed|steps?|cfg|sampler|negative\s+prompt|output\s+format)\b"
    r"|--ar\b|--stylize\b|--chaos\b|--seed\b|--q\b|16:9|4:3|宽高比|宽屏比例|分辨率|像素"
    r"|随机种子|种子|采样步数|步数|负面提示词|输出格式|模型参数|渲染参数|生成参数"
    r")"
)
ASSISTANT_META_LINE_RE = re.compile(
    r"(?i)(?:^|\s)(如果你愿意|如果需要|我可以在下一步|我可以继续|是否需要我|接下来我可以|如需我可继续)"
)
DEFAULT_INFORMATION_DENSITY = "medium"
AUTO_INFORMATION_DENSITY = "auto"
OUTLINE_INFORMATION_DENSITY_RANGES: dict[str, tuple[int, int]] = {
    "low": (1, 3),
    "medium": (3, 5),
    "high": (5, 7),
    "extra": (7, 10),
}
AUTO_OUTLINE_FALLBACK_POINTS = 3
AUTO_OUTLINE_MAX_POINTS = 10


class GenerationLogger:
    def __init__(self, run_id: str, run_dir: Path) -> None:
        self.run_id = run_id
        self.run_dir = run_dir
        self.logs_dir = run_dir / "logs"
        self.text_dir = self.logs_dir / "text"
        self.slides_dir = self.logs_dir / "slides"
        self.artifacts_dir = self.logs_dir / "artifacts"
        self.trace_path = self.logs_dir / "trace.json"
        self.progress_path = self.logs_dir / "progress.jsonl"
        self._lock = Lock()

        self.text_dir.mkdir(parents=True, exist_ok=True)
        self.slides_dir.mkdir(parents=True, exist_ok=True)
        self.artifacts_dir.mkdir(parents=True, exist_ok=True)

        now = self._now_iso()
        self.trace: dict[str, Any] = {
            "run_id": run_id,
            "run_dir": str(run_dir.resolve()),
            "created_at": now,
            "updated_at": now,
            "status": "running",
            "request": {},
            "runtime": {},
            "stages": {},
            "slides": {},
            "progress": [],
            "artifacts": {},
            "result": {},
            "error": {},
        }
        self._flush_unlocked()

    @staticmethod
    def _now_iso() -> str:
        return datetime.utcnow().isoformat(timespec="seconds") + "Z"

    def _flush_unlocked(self) -> None:
        self.trace["updated_at"] = self._now_iso()
        self.trace_path.write_text(
            json.dumps(self.trace, ensure_ascii=False, indent=2, default=str),
            encoding="utf-8",
        )

    def write_text(self, rel_path: str, text: str) -> str:
        with self._lock:
            path = self.logs_dir / rel_path
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_text(text or "", encoding="utf-8")
            return str(path.resolve())

    def write_json(self, rel_path: str, payload: Any) -> str:
        with self._lock:
            path = self.logs_dir / rel_path
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_text(json.dumps(payload, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
            return str(path.resolve())

    def write_bytes(self, rel_path: str, data: bytes) -> str:
        with self._lock:
            path = self.logs_dir / rel_path
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(data)
            return str(path.resolve())

    def set_request(self, payload: dict[str, Any]) -> None:
        with self._lock:
            self.trace["request"].update(payload)
            self._flush_unlocked()

    def set_runtime(self, payload: dict[str, Any]) -> None:
        with self._lock:
            self.trace["runtime"].update(payload)
            self._flush_unlocked()

    def set_stage(self, stage: str, payload: dict[str, Any]) -> None:
        with self._lock:
            stage_entry = self.trace["stages"].setdefault(stage, {})
            stage_entry.update(payload)
            self._flush_unlocked()

    def set_slide(self, page: int, payload: dict[str, Any]) -> None:
        with self._lock:
            slide_entry = self.trace["slides"].setdefault(str(page), {"page": page})
            slide_entry.update(payload)
            self._flush_unlocked()

    def append_slide_event(self, page: int, key: str, payload: dict[str, Any]) -> None:
        with self._lock:
            slide_entry = self.trace["slides"].setdefault(str(page), {"page": page})
            slide_entry.setdefault(key, []).append(payload)
            self._flush_unlocked()

    def set_artifact(self, name: str, value: Any) -> None:
        with self._lock:
            self.trace["artifacts"][name] = value
            self._flush_unlocked()

    def append_progress(self, payload: dict[str, Any]) -> None:
        entry = dict(payload)
        entry["timestamp"] = self._now_iso()
        with self._lock:
            with open(self.progress_path, "a", encoding="utf-8") as fp:
                fp.write(json.dumps(entry, ensure_ascii=False, default=str) + "\n")
            self.trace["progress"].append(entry)
            self._flush_unlocked()

    def finalize(self, status: str, *, result: Optional[dict[str, Any]] = None, error: Optional[dict[str, Any]] = None) -> None:
        with self._lock:
            self.trace["status"] = status
            self.trace["finished_at"] = self._now_iso()
            if result:
                self.trace["result"] = result
            if error:
                self.trace["error"] = error
            self._flush_unlocked()


class PPTImagePipeline:
    def __init__(self, settings: Settings) -> None:
        self.settings = settings
        self.output_root = Path(settings.output_root)
        self.output_root.mkdir(parents=True, exist_ok=True)

    def build_runtime_config(
        self,
        base_url: Optional[str],
        image_api_url: Optional[str],
        text_api_key: Optional[str],
        image_api_key: Optional[str],
        text_model: Optional[str],
        image_model: Optional[str],
    ) -> RuntimeConfig:
        cfg = RuntimeConfig(
            text_provider=(self.settings.text_provider or "openai").strip().lower(),
            text_base_url=(base_url or self.settings.text_base_url).strip(),
            image_provider=(self.settings.image_provider or "http").strip().lower(),
            image_base_url=(image_api_url or self.settings.image_api_url).strip(),
            text_api_key=(text_api_key or self.settings.text_api_key).strip(),
            image_api_key=(image_api_key or self.settings.resolved_image_key).strip(),
            text_model=(text_model or self.settings.text_model).strip(),
            image_model=(image_model or self.settings.image_model).strip(),
        )
        if cfg.text_provider not in {"openai", "gemini"}:
            raise ValueError("Text provider must be `openai` or `gemini`.")
        if cfg.image_provider not in {"openai", "gemini", "http"}:
            raise ValueError("Image provider must be `openai`, `gemini`, or `http`.")
        if not cfg.text_base_url:
            raise ValueError("Text model base_url cannot be empty.")
        if not cfg.image_base_url:
            raise ValueError("Image model base_url cannot be empty.")
        if not cfg.text_api_key:
            raise ValueError("Text model api_key cannot be empty.")
        if not cfg.image_api_key:
            raise ValueError("Image model api_key cannot be empty.")
        if not cfg.text_model:
            raise ValueError("Text model cannot be empty.")
        if not cfg.image_model:
            raise ValueError("Image model cannot be empty.")
        return cfg

    def run(
        self,
        user_requirement: str,
        slide_count: Optional[int],
        style_description: Optional[str],
        style_template_bytes: Optional[bytes],
        style_template_mime: Optional[str],
        runtime_cfg: RuntimeConfig,
        export_mode: str = "both",
        information_density: str = DEFAULT_INFORMATION_DENSITY,
        progress_callback: ProgressCallback = None,
    ) -> GenerateResponse:
        requirement = (user_requirement or "").strip()
        style_desc = (style_description or "").strip()
        normalized_export_mode = (export_mode or "both").strip().lower()
        normalized_information_density = self._normalize_information_density(information_density)

        if not requirement:
            raise ValueError("user_requirement cannot be empty.")
        if style_desc and style_template_bytes:
            raise ValueError("风格描述与风格模板图互斥，请二选一。")
        if normalized_export_mode not in {"images", "ppt", "both"}:
            raise ValueError("export_mode must be one of: images, ppt, both.")

        run_id = datetime.now().strftime("%Y%m%d_%H%M%S") + "_" + uuid.uuid4().hex[:8]
        run_dir = self.output_root / run_id
        run_dir.mkdir(parents=True, exist_ok=True)
        logger = GenerationLogger(run_id=run_id, run_dir=run_dir)

        style_input_meta: dict[str, Any] = {
            "type": "default",
            "description_present": bool(style_desc),
            "template_present": bool(style_template_bytes),
            "template_mime": style_template_mime or "",
        }
        if style_desc:
            style_input_meta["type"] = "description"
            style_input_meta["description_path"] = logger.write_text("artifacts/style_description.txt", style_desc)
        if style_template_bytes:
            ext = self._mime_to_extension(style_template_mime)
            template_path = logger.write_bytes(f"artifacts/style_reference{ext}", style_template_bytes)
            style_input_meta.update(
                {
                    "type": "template",
                    "template_path": template_path,
                    "template_sha256": hashlib.sha256(style_template_bytes).hexdigest(),
                    "template_size_bytes": len(style_template_bytes),
                }
            )

        logger.set_request(
            {
                "requirement": requirement,
                "requested_slide_count": slide_count,
                "information_density": normalized_information_density,
                "export_mode": normalized_export_mode,
                "style_input": style_input_meta,
            }
        )
        logger.set_runtime(
            {
                "text_provider": runtime_cfg.text_provider,
                "text_base_url": runtime_cfg.text_base_url,
                "image_provider": runtime_cfg.image_provider,
                "image_base_url": runtime_cfg.image_base_url,
                "text_model": runtime_cfg.text_model,
                "image_model": runtime_cfg.image_model,
            }
        )
        logger.set_artifact("log_dir", str(logger.logs_dir.resolve()))
        logger.set_artifact("trace_path", str(logger.trace_path.resolve()))
        logger.set_artifact("progress_log_path", str(logger.progress_path.resolve()))

        def emit(
            step: str,
            message: str,
            progress: int,
            current_slide: int = 0,
            total_slides: int = 0,
            done: bool = False,
            error: str = "",
        ) -> None:
            payload = {
                "step": step,
                "message": message,
                "progress": max(0, min(100, int(progress))),
                "current_slide": current_slide,
                "total_slides": total_slides,
                "done": done,
                "error": error,
            }
            logger.append_progress(payload)
            if not progress_callback:
                return
            progress_callback(payload)

        try:
            emit("prepare", "正在初始化模型配置...", 2)

            emit("slide_count", "正在确定PPT页数...", 8)
            resolved_slide_count = self._resolve_slide_count(
                requirement,
                slide_count,
                runtime_cfg,
                logger=logger,
            )
            logger.set_request({"resolved_slide_count": resolved_slide_count})
            emit(
                "slide_count",
                f"Slide count decided: {resolved_slide_count}",
                15,
                total_slides=resolved_slide_count,
            )

            emit("style", "Generating style prompt...", 20, total_slides=resolved_slide_count)
            style_prompt = self._generate_style_prompt(
                requirement,
                style_desc,
                style_template_bytes,
                style_template_mime,
                runtime_cfg,
                logger=logger,
            )
            emit("style", "Style prompt generated", 35, total_slides=resolved_slide_count)
            style_reference_data_url = self._image_bytes_to_data_url(style_template_bytes, style_template_mime)
            style_reference_sha256 = hashlib.sha256(style_template_bytes).hexdigest() if style_template_bytes else ""
            style_reference_mime = (style_template_mime or "image/png") if style_reference_data_url else None

            emit("outline", "Generating outline...", 38, total_slides=resolved_slide_count)
            outline = self._generate_outline(
                requirement,
                resolved_slide_count,
                normalized_information_density,
                runtime_cfg,
                logger=logger,
            )
            for slide in outline.slides:
                logger.set_slide(
                    slide.page,
                    {
                        "title": slide.title,
                        "key_points": slide.key_points,
                    },
                )
            emit("outline", "Outline generated", 50, total_slides=resolved_slide_count)

            total = len(outline.slides)
            prompt_workers = max(1, min(4, total))
            logger.set_stage(
                "prompt_generation",
                {
                    "max_workers": prompt_workers,
                    "total_slides": total,
                    "style_reference_forwarded": bool(style_reference_data_url),
                    "style_reference_mime": style_reference_mime or "",
                    "style_reference_sha256": style_reference_sha256 if style_reference_data_url else "",
                },
            )
            emit(
                "prompt_generation",
                f"开始生成每页完整 Prompt，线程数：{prompt_workers}",
                52,
                0,
                total,
            )
            completed_prompt = 0
            slide_prompts: dict[int, str] = {}
            with concurrent.futures.ThreadPoolExecutor(max_workers=prompt_workers) as prompt_pool:
                prompt_futures = {
                    prompt_pool.submit(
                        self._generate_slide_render_prompt,
                        outline.deck_title,
                        requirement,
                        slide,
                        style_prompt,
                        runtime_cfg,
                        information_density=normalized_information_density,
                        style_reference_data_url=style_reference_data_url,
                        style_reference_mime=style_reference_mime,
                        style_reference_sha256=style_reference_sha256,
                        logger=logger,
                    ): slide
                    for slide in outline.slides
                }
                try:
                    for fut in concurrent.futures.as_completed(prompt_futures):
                        slide = prompt_futures[fut]
                        slide_prompts[slide.page] = fut.result()
                        completed_prompt += 1
                        emit(
                            "prompt_generation",
                            f"Completed prompt generation {completed_prompt}/{total}",
                            52 + int((completed_prompt / max(1, total)) * 18),
                            completed_prompt,
                            total,
                        )
                except Exception:
                    for pending in prompt_futures:
                        pending.cancel()
                    raise
            emit("prompt_generation", "Slide prompts generated", 70, total, total)

            workers = max(1, min(self.settings.image_max_workers, total))
            logger.set_stage("image_generation", {"max_workers": workers, "total_slides": total})
            emit(
                "image_generation",
                f"开始并行生成图片，线程数：{workers}",
                72,
                0,
                total,
            )

            completed = 0
            results: list[SlideResult] = []
            with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as executor:
                futures = {
                    executor.submit(
                        self._render_one_slide,
                        runtime_cfg,
                        run_id,
                        run_dir,
                        slide,
                        slide_prompts.get(slide.page) or self._build_page_prompt(
                            outline.deck_title,
                            requirement,
                            slide,
                            style_prompt,
                            information_density=normalized_information_density,
                        ),
                        logger,
                    ): slide.page
                    for slide in outline.slides
                }
                try:
                    for fut in concurrent.futures.as_completed(futures):
                        results.append(fut.result())
                        completed += 1
                        emit(
                            "image_generation",
                            f"Completed image generation {completed}/{total}",
                            72 + int((completed / max(1, total)) * 20),
                            completed,
                            total,
                        )
                except Exception:
                    for pending in futures:
                        pending.cancel()
                    raise

            results.sort(key=lambda x: x.page)
            pptx_name = "generated_deck.pptx"
            pptx_path = run_dir / pptx_name
            pptx_url = ""
            resolved_pptx_path = ""
            if normalized_export_mode in {"ppt", "both"}:
                emit("packaging", "正在打包PPT...", 95, total, total)
                self._build_pptx(results, run_dir, pptx_path)
                pptx_url = f"/generated/{run_id}/{pptx_name}"
                resolved_pptx_path = str(pptx_path.resolve())
                logger.set_artifact("pptx_path", resolved_pptx_path)

            output = GenerateResponse(
                run_id=run_id,
                requirement=requirement,
                deck_title=outline.deck_title,
                style_prompt=style_prompt,
                pptx_url=pptx_url,
                pptx_path=resolved_pptx_path,
                output_dir=str(run_dir.resolve()),
                log_dir=str(logger.logs_dir.resolve()),
                trace_path=str(logger.trace_path.resolve()),
                progress_log_path=str(logger.progress_path.resolve()),
                outline=outline.slides,
                slides=results,
            )
            logger.finalize(
                "completed",
                result={
                    "deck_title": outline.deck_title,
                    "output_dir": str(run_dir.resolve()),
                    "pptx_path": resolved_pptx_path,
                    "slide_count": len(results),
                    "result_json": output.model_dump(),
                },
            )
            emit("completed", "生成完成", 100, total, total, True)
            return output
        except Exception as exc:
            logger.finalize(
                "failed",
                error={
                    "message": str(exc),
                    "traceback": traceback.format_exc(),
                },
            )
            emit("failed", f"生成失败：{exc}", 100, done=True, error=str(exc))
            raise

    def _render_one_slide(
        self,
        runtime_cfg: RuntimeConfig,
        run_id: str,
        run_dir: Path,
        slide: SlideOutline,
        prompt: str,
        logger: Optional[GenerationLogger] = None,
    ) -> SlideResult:
        file_name = f"slide_{slide.page:02d}.png"
        output_path = run_dir / file_name
        self._generate_slide_image(
            runtime_cfg,
            prompt,
            output_path,
            slide_page=slide.page,
            logger=logger,
        )
        result = SlideResult(
            page=slide.page,
            title=slide.title,
            prompt=prompt,
            image_url=f"/generated/{run_id}/{file_name}",
            image_path=str(output_path.resolve()),
        )
        if logger:
            logger.set_slide(
                slide.page,
                {
                    "image_path": result.image_path,
                    "image_url": result.image_url,
                },
            )
        return result

    def _resolve_slide_count(
        self,
        requirement: str,
        requested_slide_count: Optional[int],
        runtime_cfg: RuntimeConfig,
        logger: Optional[GenerationLogger] = None,
    ) -> int:
        if requested_slide_count is not None:
            if requested_slide_count < 1 or requested_slide_count > 20:
                raise ValueError("slide_count must be in range 1-20.")
            if logger:
                logger.set_stage(
                    "slide_count",
                    {
                        "source": "user_input",
                        "resolved_slide_count": requested_slide_count,
                    },
                )
            return requested_slide_count

        prompt = f"""
你是PPT顾问，请根据用户需求判断这份PPT最合适的页数。
输出必须是JSON，且只输出JSON：
{{
  "slide_count": 8
}}

规则：
1. 页数范围必须在 4 到 20 之间。
2. 偏简单需求建议 6-8 页，标准汇报 8-12 页，复杂复盘/方案 12-16 页。
3. 如果用户明确提到“简短”或“快速汇报”，页数偏少。
4. 如果用户明确提到“完整方案/年度复盘/多模块详解”，页数偏多。
5. 不要输出解释文本。

用户需求：
{requirement}
""".strip()

        system_prompt = "你是专业的演示文稿结构顾问。"
        stage_payload: dict[str, Any] = {}
        if logger:
            stage_payload["system_prompt_path"] = logger.write_text("text/slide_count_system.txt", system_prompt)
            stage_payload["user_prompt_path"] = logger.write_text("text/slide_count_user.txt", prompt)

        raw = chat_completion_text(
            provider=runtime_cfg.text_provider,
            base_url=runtime_cfg.text_base_url,
            api_key=runtime_cfg.text_api_key,
            model=runtime_cfg.text_model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt},
            ],
            temperature=0.2,
        )
        if logger:
            stage_payload["raw_response_path"] = logger.write_text("text/slide_count_response.txt", raw)
        try:
            count = int(self._parse_json_object(raw).get("slide_count", self.settings.default_slide_count))
        except Exception:
            nums = re.findall(r"\d+", raw or "")
            count = int(nums[0]) if nums else self.settings.default_slide_count
        resolved = max(1, min(20, count))
        if logger:
            stage_payload["source"] = "model"
            stage_payload["resolved_slide_count"] = resolved
            logger.set_stage("slide_count", stage_payload)
        return resolved

    def _generate_outline(
        self,
        requirement: str,
        slide_count: int,
        information_density: str,
        runtime_cfg: RuntimeConfig,
        logger: Optional[GenerationLogger] = None,
    ) -> OutlineResult:
        normalized_information_density = self._normalize_information_density(information_density)
        point_range = self._outline_point_range(normalized_information_density)
        density_requirement = (
            f'5. 每页 key_points 数量必须在 {point_range[0]} 到 {point_range[1]} 条之间；本次信息密度选项是 "{normalized_information_density}"。'
            if point_range
            else "5. 每页 key_points 数量由你根据页面内容自行判断，以讲清楚内容为准，不要为了整齐而机械堆砌，也不要过度稀疏。"
        )
        prompt = f"""
你是顶级演示文稿策略顾问，负责生成一套可直接用于后续出图的大纲。
你的输出必须是严格 JSON，只能输出 JSON 对象本身，不能带 markdown、解释文字或额外字段。

本次大纲以“标题 + key_points”为核心，只输出必要字段，不要写任何版式、风格、构图、镜头、视觉效果类描述。

JSON schema（严格）：
{{
  "deck_title": "string",
  "slides": [
    {{
      "page": 1,
      "title": "string",
      "key_points": ["string", "string"]
    }}
  ]
}}

硬性要求：
1. slides 数量必须严格等于 {slide_count}。
2. page 必须从 1 连续递增到 {slide_count}。
3. 所有 deck 内可见内容必须使用简体中文。
4. 每页 title 必须短、强、有区分度，不要重复表达。
{density_requirement}
6. 每条 key_point 必须把“这一点具体要讲什么”说清楚，可以是短句，也可以稍长，但不能只写一个空泛词语。
7. key_points 要便于讲述和落图，既要具体，又不要写成长段落。
8. 不要写“版式建议”“视觉重点”“画面构图”“风格说明”“图表样式”等内容。
9. 如果用户需求很宽泛，你要主动补足合理结构和必要子主题，但不要反问用户。
10. 如果用户需求很窄，你要在不跑题的前提下补出能撑起完整汇报的内容层次。
11. 避免不同页面承担同样作用；整套内容要形成清晰叙事推进。

质量标准：
- 每页都要回答“为什么要有这一页”。
- 整体叙事建议遵循：开场 -> 背景/问题 -> 分析/拆解 -> 方案/路径 -> 收益/落地 -> 收尾。
- 每页 key_points 之间要有内部层次，不要只是同义改写。
- 内容应符合真实汇报场景，适合管理层/业务/项目沟通使用。

用户需求：
{requirement}
""".strip()

        system_prompt = (
            "你生成严格 JSON 格式的 PPT 大纲。"
            "只输出 deck_title 和 slides.page/title/key_points 这几个字段。"
        )
        stage_payload: dict[str, Any] = {
            "information_density": normalized_information_density,
        }
        if point_range:
            stage_payload["key_points_min"] = point_range[0]
            stage_payload["key_points_max"] = point_range[1]
        if logger:
            stage_payload["system_prompt_path"] = logger.write_text("text/outline_system.txt", system_prompt)
            stage_payload["user_prompt_path"] = logger.write_text("text/outline_user.txt", prompt)

        raw_response = chat_completion_text(
            provider=runtime_cfg.text_provider,
            base_url=runtime_cfg.text_base_url,
            api_key=runtime_cfg.text_api_key,
            model=runtime_cfg.text_model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt},
            ],
            temperature=0.45,
        )
        if logger:
            stage_payload["raw_response_path"] = logger.write_text("text/outline_response.txt", raw_response)
        parsed = self._parse_json_object(raw_response)
        slides_data = parsed.get("slides")
        if not isinstance(slides_data, list):
            raise ValueError("Outline parse failed: slides is invalid.")

        slides: list[SlideOutline] = []
        for i, item in enumerate(slides_data, start=1):
            points = self._normalize_outline_key_points(item.get("key_points", []), normalized_information_density)
            slides.append(
                SlideOutline(
                    page=int(item.get("page", i)),
                    title=str(item.get("title", f"第{i}页")),
                    key_points=points,
                )
            )

        if len(slides) != slide_count:
            slides = self._normalize_slides(slides, slide_count, normalized_information_density)
        outline = OutlineResult(deck_title=str(parsed.get("deck_title", "自动生成大纲")), slides=slides)
        if logger:
            outline_json_path = logger.write_json(
                "text/outline_parsed.json",
                outline.model_dump(),
            )
            stage_payload["deck_title"] = outline.deck_title
            stage_payload["outline_json_path"] = outline_json_path
            stage_payload["slide_count"] = len(outline.slides)
            logger.set_stage("outline", stage_payload)
            logger.set_artifact("outline_json_path", outline_json_path)
        return outline

    def _normalize_slides(
        self,
        slides: list[SlideOutline],
        slide_count: int,
        information_density: str = DEFAULT_INFORMATION_DENSITY,
    ) -> list[SlideOutline]:
        normalized_information_density = self._normalize_information_density(information_density)
        out: list[SlideOutline] = []
        for i in range(1, slide_count + 1):
            if i <= len(slides):
                raw = slides[i - 1]
                out.append(
                    SlideOutline(
                        page=i,
                        title=raw.title or f"第{i}页",
                        key_points=self._normalize_outline_key_points(raw.key_points, normalized_information_density),
                    )
                )
            else:
                out.append(
                    SlideOutline(
                        page=i,
                        title=f"第{i}页",
                        key_points=self._fallback_outline_key_points(i, normalized_information_density),
                    )
                )
        return out

    @staticmethod
    def _normalize_information_density(information_density: str) -> str:
        normalized = (information_density or DEFAULT_INFORMATION_DENSITY).strip().lower()
        if not normalized:
            normalized = DEFAULT_INFORMATION_DENSITY
        if normalized not in {AUTO_INFORMATION_DENSITY, *OUTLINE_INFORMATION_DENSITY_RANGES.keys()}:
            raise ValueError(
                "information_density must be one of: auto, low, medium, high, extra."
            )
        return normalized

    @staticmethod
    def _outline_point_range(information_density: str) -> Optional[tuple[int, int]]:
        normalized = PPTImagePipeline._normalize_information_density(information_density)
        if normalized == AUTO_INFORMATION_DENSITY:
            return None
        return OUTLINE_INFORMATION_DENSITY_RANGES[normalized]

    def _fallback_outline_key_points(self, page: int, information_density: str) -> list[str]:
        point_range = self._outline_point_range(information_density)
        min_points = point_range[0] if point_range else AUTO_OUTLINE_FALLBACK_POINTS
        prefix = f"第 {page} 页" if page > 0 else "当前页"
        return [
            f"补充{prefix}需要讲清楚的关键信息 {idx}。"
            for idx in range(1, min_points + 1)
        ]

    @staticmethod
    def _build_information_density_guidance(information_density: str) -> str:
        normalized = PPTImagePipeline._normalize_information_density(information_density)
        if normalized == AUTO_INFORMATION_DENSITY:
            return ""
        point_range = PPTImagePipeline._outline_point_range(normalized)
        if not point_range:
            return ""
        min_points, max_points = point_range
        guidance_map = {
            "low": [
                f"本页属于低信息密度页面，建议只承载 {min_points}-{max_points} 个主要信息单元。",
                "优先使用更大的模块、更明确的主次关系和更充足的留白，不要把页面塞满。",
                "可以让单个模块里的说明稍微完整一些，但整体模块数量要克制，避免出现拥挤拼贴感。",
            ],
            "medium": [
                f"本页属于中等信息密度页面，建议承载 {min_points}-{max_points} 个主要信息单元。",
                "在信息完整度和留白之间保持平衡，适合一主多辅或多分区但不过载的结构。",
                "模块数量适中，既要体现内容层次，也要保留足够呼吸感，避免堆砌成高密度看板。",
            ],
            "high": [
                f"本页属于高信息密度页面，建议承载 {min_points}-{max_points} 个主要信息单元。",
                "可以使用更多分区、卡片、流程节点或对比模块，让页面信息更丰富，但必须保持稳定栅格、清晰分层和可扫读性。",
                "允许间距更紧凑，但不能牺牲标题锚点、主信息入口和阅读路径。",
            ],
            "extra": [
                f"本页属于超高信息密度页面，建议承载 {min_points}-{max_points} 个主要信息单元。",
                "应主动考虑矩阵、多栏、看板、复杂流程、架构图或仪表盘式组织方式，把大量信息装进统一骨架。",
                "允许更紧凑的排布和更高模块数量，但必须通过强层级、稳定对齐和分区节奏保证页面仍然清晰可读。",
            ],
        }
        return PPTImagePipeline._format_guidance_lines(guidance_map[normalized], limit=4)

    def _ensure_prompt_density_guidance(self, prompt: str, information_density: str) -> str:
        normalized_prompt = (prompt or "").strip()
        if not normalized_prompt:
            normalized_prompt = "完整宽屏中文 PPT 单页。"
        if self._normalize_information_density(information_density) == AUTO_INFORMATION_DENSITY:
            return normalized_prompt
        if "本页信息密度辅助约束：" in normalized_prompt or "信息密度控制：" in normalized_prompt:
            return normalized_prompt
        density_guidance = self._build_information_density_guidance(information_density)
        return f"{normalized_prompt}\n\n信息密度控制：\n{density_guidance}".strip()

    def _normalize_outline_key_points(self, raw_points: Any, information_density: str) -> list[str]:
        normalized_density = self._normalize_information_density(information_density)
        point_range = self._outline_point_range(normalized_density)
        if point_range:
            min_points, max_points = point_range
        else:
            min_points, max_points = 1, AUTO_OUTLINE_MAX_POINTS
        if isinstance(raw_points, list):
            candidates = raw_points
        elif raw_points is None:
            candidates = []
        else:
            candidates = [raw_points]

        normalized_points: list[str] = []
        seen: set[str] = set()
        for raw_point in candidates:
            point = re.sub(r"\s+", " ", str(raw_point or "").strip())
            point = re.sub(r"^[-*•·●]\s*", "", point)
            point = re.sub(r"^\d+[\.、]\s*", "", point)
            point = point.strip()
            if not point:
                continue
            if not re.search(r"[。！？；;]$", point):
                point += "。"
            lowered = point.lower()
            if lowered in seen:
                continue
            seen.add(lowered)
            normalized_points.append(point)
            if len(normalized_points) >= max_points:
                break

        if normalized_density == AUTO_INFORMATION_DENSITY:
            if normalized_points:
                return normalized_points[:AUTO_OUTLINE_MAX_POINTS]
            return self._fallback_outline_key_points(0, normalized_density)
        if len(normalized_points) < min_points:
            normalized_points.extend(
                self._fallback_outline_key_points(0, normalized_density)[: min_points - len(normalized_points)]
            )
        return normalized_points[:max_points]

    def _generate_style_prompt(
        self,
        requirement: str,
        style_description: str,
        style_template_bytes: Optional[bytes],
        style_template_mime: Optional[str],
        runtime_cfg: RuntimeConfig,
        logger: Optional[GenerationLogger] = None,
    ) -> str:
        fallback = self._build_default_style_prompt(requirement)
        if style_description:
            user_prompt = f"""
你是顶级视觉策略顾问。现在请把用户提供的“风格描述”扩写为一段超详细、可复用、可直接用于整套 PPT 统一视觉的风格说明。

要求：
1. 不设字数上限，请尽可能详细、具体、充分。
2. 重点描写非内容部分的风格样式，而不是页面具体文案内容。
3. 需要详细到足以让另一个模型仅凭这段风格说明，就能尽量复现出风格高度一致、气质统一的图片。
4. 请重点展开描述：整体气质、版式骨架、栅格与留白、背景系统、色彩关系、标题与正文的视觉关系、字体气质、卡片/面板/分区样式、图标与插画语言、图表外观、线条与边框、阴影与质感、层次与景深、装饰元素、节奏感、页面动线、视觉锚点、统一性规则、禁用项。
5. 明确说明哪些是整套 PPT 中必须始终保持一致的视觉特征，哪些只能做有限变化。
6. 不要把重点放在某一页的业务内容上，而要抽取跨页复用的视觉规则。
7. 不要输出 PPT 生成参数或技术指令，例如宽高比、分辨率、像素、seed、steps、CFG、输出格式、模型参数等。
8. 输出语言必须是中文。
9. 输出格式不做限制，不必按固定栏目组织，你可以用最适合复现风格的方式自由详细描述。
10. 不要解释分析过程，直接给出最终风格描述。

用户风格描述：
{style_description}
""".strip()
            system_prompt = "你负责把风格描述扩写为超详细、强可复现、可跨页复用的视觉风格说明。"
            stage_payload: dict[str, Any] = {"source": "description"}
            if logger:
                stage_payload["system_prompt_path"] = logger.write_text("text/style_system.txt", system_prompt)
                stage_payload["user_prompt_path"] = logger.write_text("text/style_user.txt", user_prompt)
                stage_payload["request_payload_path"] = logger.write_json(
                    "text/style_request.json",
                    {
                        "mode": "description",
                        "model": runtime_cfg.text_model,
                        "temperature": 0.2,
                        "messages": [
                            {"role": "system", "content_type": "text", "content_path": stage_payload["system_prompt_path"]},
                            {"role": "user", "content_type": "text", "content_path": stage_payload["user_prompt_path"]},
                        ],
                        "image_attachment_present": False,
                    },
                )
            try:
                raw_response = chat_completion_text(
                    provider=runtime_cfg.text_provider,
                    base_url=runtime_cfg.text_base_url,
                    api_key=runtime_cfg.text_api_key,
                    model=runtime_cfg.text_model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt},
                    ],
                    temperature=0.2,
                )
                raw_response = raw_response.strip()
                final_prompt = self._finalize_style_prompt(raw_response or fallback, fallback)
                if logger:
                    stage_payload["raw_response_path"] = logger.write_text("text/style_response.txt", raw_response)
                    final_prompt_path = logger.write_text("artifacts/style_prompt.txt", final_prompt)
                    stage_payload["final_prompt_path"] = final_prompt_path
                    stage_payload["fallback_used"] = not bool(raw_response)
                    logger.set_stage("style", stage_payload)
                    logger.set_artifact("style_prompt_path", final_prompt_path)
                return final_prompt
            except Exception as exc:
                final_prompt = self._finalize_style_prompt(fallback, fallback)
                if logger:
                    final_prompt_path = logger.write_text("artifacts/style_prompt.txt", final_prompt)
                    stage_payload["final_prompt_path"] = final_prompt_path
                    stage_payload["fallback_used"] = True
                    stage_payload["error"] = str(exc)
                    logger.set_stage("style", stage_payload)
                    logger.set_artifact("style_prompt_path", final_prompt_path)
                return final_prompt

        if style_template_bytes:
            mime = style_template_mime or "image/png"
            data_url = self._image_bytes_to_data_url(style_template_bytes, mime) or ""
            user_prompt = f"""
你是高级视觉分析师。请分析这张 PPT 风格模板图，并输出一段超详细、可复用、可直接用于整套 PPT 统一视觉的中文风格说明。

要求：
1. 不设字数上限，请尽可能详细、具体、充分。
2. 重点识别和描述这张图里非内容部分的风格样式，而不是照着图片里已有文字内容做摘要。
3. 需要把视觉风格拆解到足够细，让另一个模型仅凭你的描述，就能尽量复现出风格高度一致、气质统一的图片。
4. 请详细描述：整体气质、背景处理、配色关系、明暗结构、留白策略、版式骨架、标题区和内容区的样式关系、卡片与容器、图标与插画语言、图表样式、边框与线条、阴影与质感、装饰元素、页面动线、视觉焦点组织方式、统一性规则、禁用项。
5. 请特别指出哪些视觉特征是这套风格最不可丢失的核心特征，哪些变化会破坏一致性。
6. 你要总结的是整套 PPT 可以复用的视觉 DNA，而不是对单页内容主题做解释。
7. 不要输出 PPT 生成参数或技术指令，例如宽高比、分辨率、像素、seed、steps、CFG、输出格式、模型参数等。
8. 输出语言必须是中文。
9. 请尽可能写得更长、更细，不要遗漏任何会影响风格复现的重要细节。
10. 不要解释你的思考过程，直接给出最终风格描述。
""".strip()
            system_prompt = "你输出超详细、强可复现、可跨页复用的视觉风格说明。"
            stage_payload = {"source": "template", "template_mime": mime}
            if logger:
                stage_payload["system_prompt_path"] = logger.write_text("text/style_system.txt", system_prompt)
                stage_payload["user_prompt_path"] = logger.write_text("text/style_user.txt", user_prompt)
                stage_payload["request_payload_path"] = logger.write_json(
                    "text/style_request.json",
                    {
                        "mode": "template",
                        "model": runtime_cfg.text_model,
                        "temperature": 0.2,
                        "messages": [
                            {"role": "system", "content_type": "text", "content_path": stage_payload["system_prompt_path"]},
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text_path": stage_payload["user_prompt_path"]},
                                    {
                                        "type": "image_url",
                                        "attached": True,
                                        "mime": mime,
                                        "data_url_prefix": data_url[:64],
                                        "data_url_length": len(data_url),
                                        "image_sha256": hashlib.sha256(style_template_bytes).hexdigest(),
                                    },
                                ],
                            },
                        ],
                        "image_attachment_present": True,
                    },
                )
            try:
                raw_response = chat_completion_text(
                    provider=runtime_cfg.text_provider,
                    base_url=runtime_cfg.text_base_url,
                    api_key=runtime_cfg.text_api_key,
                    model=runtime_cfg.text_model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": user_prompt},
                                {"type": "image_url", "image_url": {"url": data_url}},
                            ],
                        },
                    ],
                    temperature=0.2,
                )
                raw_response = raw_response.strip()
                final_prompt = self._finalize_style_prompt(raw_response or fallback, fallback)
                if logger:
                    stage_payload["raw_response_path"] = logger.write_text("text/style_response.txt", raw_response)
                    final_prompt_path = logger.write_text("artifacts/style_prompt.txt", final_prompt)
                    stage_payload["final_prompt_path"] = final_prompt_path
                    stage_payload["fallback_used"] = not bool(raw_response)
                    logger.set_stage("style", stage_payload)
                    logger.set_artifact("style_prompt_path", final_prompt_path)
                return final_prompt
            except Exception as exc:
                final_prompt = self._finalize_style_prompt(fallback, fallback)
                if logger:
                    final_prompt_path = logger.write_text("artifacts/style_prompt.txt", final_prompt)
                    stage_payload["final_prompt_path"] = final_prompt_path
                    stage_payload["fallback_used"] = True
                    stage_payload["error"] = str(exc)
                    logger.set_stage("style", stage_payload)
                    logger.set_artifact("style_prompt_path", final_prompt_path)
                return final_prompt

        return self._infer_style_prompt_from_requirement(
            requirement=requirement,
            runtime_cfg=runtime_cfg,
            fallback=fallback,
            logger=logger,
        )

    def _infer_style_prompt_from_requirement(
        self,
        requirement: str,
        runtime_cfg: RuntimeConfig,
        fallback: str,
        logger: Optional[GenerationLogger] = None,
    ) -> str:
        user_prompt = f"""
你是顶级演示视觉策略顾问。当前没有用户提供的风格描述，也没有风格参考图。
请你根据“用户需求本身”反推一套最合适的整套 PPT 视觉风格 DNA，并输出一段超详细、可复用、可直接用于后续所有页面统一出图的中文风格说明。

你的任务不是复述业务内容，而是先理解需求中的行业属性、受众身份、汇报场景、表达目标、信息密度、说服方式，再为这份 PPT 设计一套匹配的视觉系统。

要求：
1. 输出必须是中文。
2. 不设字数上限，请尽可能长、尽可能完整、尽可能具体、尽可能细致。
3. 必须从需求里主动判断并明确风格方向，例如更偏董事会汇报、科研汇报、产业方案、政企汇报、技术架构说明、品牌提案、教育培训、医疗说明等，但不要把它写成一句抽象判断，要把这种判断落实为视觉规则。
4. 重点描述整套 PPT 跨页复用的非内容风格规则，而不是某一页的具体业务文案。
5. 请系统展开描述：整体气质、视觉定位、受众感知、背景系统、配色架构、标题系统、正文字体气质、栅格与留白、模块容器、卡片语言、图标与插画体系、图表样式、边框与线条、阴影与材质、信息密度、视觉焦点组织、页面动线、统一性规则、禁用项。
6. 必须明确哪些视觉特征必须始终保持一致，哪些地方允许有限变化。
7. 如果需求显得偏正式、偏理性、偏技术、偏政企、偏管理层，请避免自动写成泛化、花哨、互联网营销感的风格；如果需求本身偏创意、偏品牌、偏年轻化，也要忠实体现，不要强行改成传统商务蓝灰模板。
8. 风格说明必须足够细，让另一个模型仅凭这段文字就能稳定生成整套风格一致的 PPT 页面。
9. 不要输出 PPT 生成参数或技术指令，例如宽高比、分辨率、像素、seed、steps、CFG、输出格式、模型参数等。
10. 不要解释你的思考过程，直接给出最终风格说明。

用户需求：
{requirement}
""".strip()
        system_prompt = "你负责根据用户需求反推最合适的 PPT 视觉风格 DNA，并输出超详细、强可复现、可跨页复用的风格说明。"
        stage_payload: dict[str, Any] = {"source": "requirement_inference"}
        if logger:
            stage_payload["system_prompt_path"] = logger.write_text("text/style_system.txt", system_prompt)
            stage_payload["user_prompt_path"] = logger.write_text("text/style_user.txt", user_prompt)
            stage_payload["request_payload_path"] = logger.write_json(
                "text/style_request.json",
                {
                    "mode": "requirement_inference",
                    "model": runtime_cfg.text_model,
                    "temperature": 0.25,
                    "messages": [
                        {"role": "system", "content_type": "text", "content_path": stage_payload["system_prompt_path"]},
                        {"role": "user", "content_type": "text", "content_path": stage_payload["user_prompt_path"]},
                    ],
                    "image_attachment_present": False,
                },
            )
        try:
            raw_response = chat_completion_text(
                provider=runtime_cfg.text_provider,
                base_url=runtime_cfg.text_base_url,
                api_key=runtime_cfg.text_api_key,
                model=runtime_cfg.text_model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
                temperature=0.25,
            )
            raw_response = raw_response.strip()
            final_prompt = self._finalize_style_prompt(raw_response or fallback, fallback)
            if logger:
                stage_payload["raw_response_path"] = logger.write_text("text/style_response.txt", raw_response)
                final_prompt_path = logger.write_text("artifacts/style_prompt.txt", final_prompt)
                stage_payload["final_prompt_path"] = final_prompt_path
                stage_payload["fallback_used"] = not bool(raw_response)
                logger.set_stage("style", stage_payload)
                logger.set_artifact("style_prompt_path", final_prompt_path)
            return final_prompt
        except Exception as exc:
            final_prompt = self._finalize_style_prompt(fallback, fallback)
            if logger:
                final_prompt_path = logger.write_text("artifacts/style_prompt.txt", final_prompt)
                stage_payload["final_prompt_path"] = final_prompt_path
                stage_payload["fallback_used"] = True
                stage_payload["error"] = str(exc)
                logger.set_stage("style", stage_payload)
                logger.set_artifact("style_prompt_path", final_prompt_path)
            return final_prompt

    def _build_default_style_prompt(self, requirement: str) -> str:
        return f"""
[Style Blueprint | Long Form | Highly Reproducible]
This project needs a complete, reusable, high-resolution, image-first PPT visual language.
The style must feel premium, minimal, highly controlled, strongly designed, and narrative-driven, suitable for formal Chinese business storytelling.

Project context:
{requirement}

1) Macro Creative Direction
Build a modern strategic-report visual system with strong hierarchy, controlled visual rhythm, and persuasive framing.
Every page should look intentional, not generic, with clear focal hierarchy and content zones.
Design temperament should balance rational analysis, simplified visual expression, and compelling storytelling.
The overall feel should be calm, precise, premium, restrained, and cohesive.
The style should not rely on novelty from page to page. Instead, it should create the impression that all slides come from the same tightly managed art direction system.

2) Color Architecture
Define a disciplined color architecture:
- Primary palette for authority and brand anchoring.
- Secondary palette for supportive areas and module separation.
- Accent palette for key numbers, breakthroughs, risks, and calls-to-action.
- Neutral ramp for structure, backgrounds, and text contrast.
Apply contrast with purpose, not decoration.
Ensure consistency across pages while allowing controlled variation by section.
Background tones, panel tones, dividers, highlights, and text contrast should always stay within one narrow family of related values.
Avoid random per-slide recoloring.

3) Typography and Readability
Use high-legibility Chinese-friendly typographic behavior:
- Strong title style with high contrast against background.
- Subtitles to create sectional pacing.
- Body text regions must remain readable at presentation distance.
- Numeric emphasis style for data highlights.
Maintain stable spacing scale so page rhythm is coherent.
Typography should feel consistent in weight distribution, density, alignment logic, and title-to-body proportion across every slide.
Headings, labels, captions, annotations, and emphasized figures should feel like one family rather than independently designed elements.

4) Layout and Grid Logic
Use a robust widescreen composition strategy:
- Keep clear safe margins.
- Favor modular grids.
- Reserve predictable anchors for title, core insight, and evidence blocks.
- Balance dense information regions with whitespace buffers.
- Prevent visual crowding near edges.
Let the page scaffold feel repeatable from slide to slide: similar margin behavior, similar internal padding logic, similar content block proportions, and similar alignment rhythm.
The viewer should sense the same layout DNA even when slide structures differ.

5) Visual Components
Define repeatable components:
- Insight cards
- KPI highlights
- Comparison matrices
- Timeline modules
- Process flow blocks
- Risk/opportunity chips
- Evidence panels with concise annotation zones
Each component should share one visual DNA but adapt to page intent.
Card corners, border behavior, fill treatment, divider treatment, corner radii, and surface layering should remain highly consistent across the deck.

6) Illustration and Icon Style
Use one consistent icon/illustration family:
- Controlled detail level
- Professional and clean contour
- Avoid childish or playful mismatch
- Harmonize line thickness and fill behavior
- Keep symbolism business-relevant and easy to decode
If illustration is used, it must belong to one stable family in rendering method, abstraction level, edge quality, and color treatment.
Do not mix realistic, flat, pseudo-3D, sketch-like, and glossy illustration languages in one deck.

7) Data Storytelling Expression
Charts should read as strategic evidence, not decorative graphics:
- Strong data-ink ratio
- Clear labeling hierarchy
- Accent only for key insights
- Secondary values visually subdued
- Annotation language concise and persuasive
Axes, legends, labels, callouts, data markers, bars, lines, and comparison modules should share a consistent chart styling system.

8) Depth, Lighting, and Texture
Keep depth subtle and intentional:
- Soft layering and tonal separation
- Controlled shadows
- Occasional gradient transitions
- Minimal texture noise
- Avoid over-stylized visual effects that reduce readability
- Keep decoration restrained and minimal rather than flashy
Shadow softness, layer stacking, glassiness or matte feeling, glow usage, blur usage, and gradient transitions should be managed as a stable system rather than varied casually.

9) Page-by-Page Consistency Rules
Maintain continuity across all slides:
- Common baseline spacing units
- Repeatable heading rhythm
- Stable color role mapping
- Predictable module behavior
- Consistent visual punctuation and separators
The following should remain especially stable:
- background treatment
- title styling
- panel/chip/card styling
- icon family
- chart appearance
- edge treatment and corner language
- highlight color behavior
- shadow and depth behavior
- density and whitespace ratio

The following may vary only slightly:
- section-level accent emphasis
- page-specific composition shape
- amount of diagrammatic structure
- intensity of decorative support graphics

The following should never drift:
- overall visual temperament
- rendering language
- contrast logic
- typography hierarchy logic
- surface material feel
- visual cleanliness standard

10) Chinese Business Context Constraints
All visible text should be in Simplified Chinese.
Tone should be executive-friendly, objective, and insight-heavy.
Avoid slang, meme aesthetics, and excessive ornament.
Prioritize trust, clarity, and strategic confidence.

11) Consistency Safeguards
Do not impose a preset taste bias that overrides the intended style.
If the chosen style naturally contains stronger decoration, illustration, collage, contrast, density, or expressive visual treatment, preserve it faithfully.
Only avoid visual decisions that conflict with the established style system, reduce readability below acceptable presentation standards, or break deck-level continuity without justification from the style itself.
Prevent accidental drift in rendering language, contrast logic, typography hierarchy, component behavior, and material feeling unless the style blueprint itself clearly calls for those shifts.

12) Style-Only Constraint
Only describe visual language and reusable layout rules.
Do not include PPT generation parameters, aspect-ratio settings, resolution, pixels, seeds, sampling steps, CFG, output-format instructions, or model-setting content.

13) Reusable Generation Clause
This style blueprint is meant to be reused across many slides.
Any generated page should preserve this visual DNA while adapting to that page's specific narrative purpose.
When conflict occurs, preserve readability, hierarchy, and business clarity first.
""".strip()

    @staticmethod
    def _dedupe_lines(lines: list[str]) -> list[str]:
        unique: list[str] = []
        seen: set[str] = set()
        for line in lines:
            cleaned = re.sub(r"\s+", " ", (line or "").strip())
            if not cleaned:
                continue
            key = cleaned.lower()
            if key in seen:
                continue
            seen.add(key)
            unique.append(cleaned)
        return unique

    @staticmethod
    def _format_guidance_lines(lines: list[str], limit: int) -> str:
        formatted: list[str] = []
        for line in PPTImagePipeline._dedupe_lines(lines)[:limit]:
            cleaned = line.rstrip()
            if not re.search(r"[。！？；;：:，,\.\!]$", cleaned):
                cleaned += "。"
            formatted.append(f"- {cleaned}")
        return "\n".join(formatted).strip()

    def _finalize_style_prompt(self, raw_text: str, fallback_text: str) -> str:
        cleaned = self._sanitize_style_prompt(raw_text)
        if not cleaned:
            cleaned = self._sanitize_style_prompt(fallback_text)
        if not cleaned:
            cleaned = STYLE_PROMPT_GUARD
        return cleaned

    def _sanitize_style_prompt(self, raw_text: str) -> str:
        lines: list[str] = []
        for line in re.split(r"\r?\n+", raw_text or ""):
            original = line.strip()
            normalized = re.sub(r"\s+", " ", original).strip()
            if not normalized:
                continue
            if ASSISTANT_META_LINE_RE.search(normalized):
                break
            if STYLE_PROMPT_PARAMETER_RE.search(normalized) and not STYLE_PROMPT_PROHIBITION_RE.search(normalized):
                continue
            lines.append(original)
        return "\n".join(lines).strip()

    def _generate_slide_render_prompt(
        self,
        deck_title: str,
        requirement: str,
        slide: SlideOutline,
        style_prompt: str,
        runtime_cfg: RuntimeConfig,
        information_density: str = DEFAULT_INFORMATION_DENSITY,
        style_reference_data_url: Optional[str] = None,
        style_reference_mime: Optional[str] = None,
        style_reference_sha256: str = "",
        logger: Optional[GenerationLogger] = None,
    ) -> str:
        points = "\n".join(f"- {x}" for x in slide.key_points) or "- TBD point"
        has_style_reference = bool(style_reference_data_url)
        normalized_density = self._normalize_information_density(information_density)
        density_guidance = self._build_information_density_guidance(normalized_density)
        density_section = (
            f"\n本页信息密度参考：\n{density_guidance}\n"
            if density_guidance
            else ""
        )
        reference_guidance = """
风格参考图约束：
- 已提供风格参考图。它是整套 PPT 非内容视觉风格的最高优先级真值与母版。
- 你生成的详细 prompt 必须让最终页面继续属于这张参考图的同一视觉家族：相同的标题系统、色彩关系、分区方式、模块框语言、边框/线条习惯、图解组织方式、图标语法、材质气质、信息密度与版式节奏。
- 风格文字说明只是对参考图的补充解释，不得覆盖、弱化或重写参考图本身。
- 不允许为了“更简洁、更现代、更像通用商务模板”而改造成另一套视觉系统。
- 必须尽量复刻参考图的非内容特征，只替换为当前页需要承载的主题内容与信息结构。
- 如果参考图本身偏高密度、工程图解、科研汇报、流程框图、模块矩阵、强分区布局，就必须保留这种结构复杂度和图解气质，不能自动扁平化成简单卡片流程图。
""".strip() if has_style_reference else """
风格参考图约束：
- 本次没有附带额外风格参考图，请严格依据风格 DNA 生成当前页详细 prompt。
""".strip()
        user_prompt = f"""
你是世界级演示视觉总监，也是 AI 绘图提示词专家。
请基于整套 PPT 的风格 DNA、版式骨架和当前页大纲，生成“这一页用于出图的完整 prompt”。

硬性要求：
1. 只输出最终 prompt 正文，不要解释，不要加标题。
2. 输出必须尽可能长、尽可能完整、尽可能具体、尽可能细致，宁可更详细也不要因为概括而漏掉细节。
3. 必须把页面里每一个重要可见元素都描述出来，包括但不限于：版头、标题区、副标题或说明文字、主视觉区、内容分区、卡片/模块、图表/信息图、标签、箭头、连接线、边框、图标、装饰线、底部辅助区、页码区、重点强调色位置、阴影与质感处理。
4. 必须明确页面结构、信息层级、主视觉焦点、内容区块、图表/信息图表达、标题与正文的呈现方式、留白、材质、光影和视觉节奏。
5. 必须把各个模块之间的相对位置、大小关系、对齐关系、主次层级、视觉流向写清楚，避免只给抽象概述。
6. 对每个关键区块，尽量写清楚它承载什么信息、长什么样、如何与周围元素连接，而不是只说“放一个模块”。
7. 必须严格继承给定风格 DNA，但要转化为当前页可执行的画面描述，不要机械照抄整段风格文本。
8. 所有需要出现在页内的文字都必须是简体中文，并且内容要与当前页主题一致。
9. 结果必须像完整宽屏 PPT 单页，能够完整承载页面结构与信息层级，而不是只呈现局部元素或缺少页面整体骨架。
10. 不要输出分辨率、像素、seed、steps、CFG、模型参数、输出格式等技术参数。
11. 如果当前页包含数据、流程、对比、时间线或分层结构，要在 prompt 中明确相应的信息图表达方式。
12. 必须保留当前页标题和关键要点，不得遗漏核心内容。
13. 必须把“整套 PPT 的风格一致性”作为硬约束，优先级高于单页新鲜感和局部创意发挥。
14. 背景处理、配色关系、卡片样式、边框语言、阴影质感、图标/插画体系、图表外观、排版气质必须与其他页保持同一视觉家族，不允许风格漂移。
15. 允许内容布局按当前页主题变化，但不允许视觉风格体系变化。
16. 如果存在取舍，优先保证与整套风格 DNA 的一致性，再考虑单页表现力。
17. 如果附带了风格参考图，必须把参考图视为最高优先级依据；风格 DNA 只是帮助你把参考图中的视觉规律解释得更清楚。
18. 如果参考图呈现的是某种特定的图解密度、模块结构、标题带样式、强调色策略、边框与箭头语言，你必须尽量在 prompt 中把这些非内容特征落地到当前页；如果参考图本身就带有更强装饰性、插画感、漫画感、拼贴感、极简感或其他鲜明表达，也应在保持一致性的前提下忠实保留。
19. 不要用预设审美偏好去修正参考图或风格 DNA；只需要避免偏离、误读或扁平化成与原始风格不一致的另一套视觉体系。
20. 输出时要像在给图像模型做逐层施工说明，尽量覆盖从整体布局到局部样式的全部关键细节，不要遗漏任何会影响最终画面一致性和信息完整性的点。

整套上下文：
Deck title: {deck_title}
User requirement: {requirement}

当前页信息：
Slide index: {slide.page}
Slide title: {slide.title}
Key points:
{points}
{density_section}
{reference_guidance}

风格 DNA：
{style_prompt}
""".strip()
        fallback = self._build_page_prompt(
            deck_title,
            requirement,
            slide,
            style_prompt,
            information_density=information_density,
        )
        system_prompt = (
            "你生成可直接用于 PPT 单页出图的完整提示词，只输出提示词正文。"
            "提示词必须尽可能详细、尽可能完整，宁可更长也不要遗漏任何结构、文案、样式和视觉细节。"
            "如果用户消息中附带风格参考图，你必须把它当作最高优先级视觉真值，"
            "不要用预设审美偏好去改写它，只需忠实延展其视觉系统。"
        )
        slide_dir = f"slides/slide_{slide.page:02d}"
        request_payload: dict[str, Any] = {
            "provider": runtime_cfg.text_provider,
            "model": runtime_cfg.text_model,
            "temperature": 0.1,
            "information_density": self._normalize_information_density(information_density),
            "image_attachment_present": has_style_reference,
            "messages": [
                {"role": "system", "content_type": "text"},
                {
                    "role": "user",
                    "content": [
                        {"type": "text"},
                    ],
                },
            ],
        }
        if logger:
            system_prompt_path = logger.write_text(f"{slide_dir}/prompt_system.txt", system_prompt)
            user_prompt_path = logger.write_text(f"{slide_dir}/prompt_user.txt", user_prompt)
            request_payload["messages"][0]["content_path"] = system_prompt_path
            request_payload["messages"][1]["content"][0]["text_path"] = user_prompt_path
            if has_style_reference:
                request_payload["messages"][1]["content"].append(
                    {
                        "type": "image_url",
                        "attached": True,
                        "mime": style_reference_mime or "image/png",
                        "data_url_prefix": (style_reference_data_url or "")[:64],
                        "data_url_length": len(style_reference_data_url or ""),
                        "image_sha256": style_reference_sha256,
                    }
                )
            logger.set_slide(
                slide.page,
                {
                    "prompt_system_path": system_prompt_path,
                    "prompt_user_path": user_prompt_path,
                    "prompt_fallback_path": logger.write_text(f"{slide_dir}/prompt_fallback.txt", fallback),
                    "prompt_request_path": logger.write_json(f"{slide_dir}/prompt_request.json", request_payload),
                    "prompt_reference_image_attached": has_style_reference,
                    "prompt_reference_image_sha256": style_reference_sha256 if has_style_reference else "",
                },
            )
        try:
            user_content: Any = user_prompt
            if has_style_reference:
                user_content = [
                    {"type": "text", "text": user_prompt},
                    {"type": "image_url", "image_url": {"url": style_reference_data_url}},
                ]
            raw_output = chat_completion_text(
                provider=runtime_cfg.text_provider,
                base_url=runtime_cfg.text_base_url,
                api_key=runtime_cfg.text_api_key,
                model=runtime_cfg.text_model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_content},
                ],
                temperature=0.1,
            ).strip()
            final_prompt = self._ensure_prompt_density_guidance(
                self._normalize_long_text_prompt(raw_output or fallback),
                information_density,
            )
            if logger:
                payload = {
                    "prompt_response_path": logger.write_text(f"{slide_dir}/prompt_response.txt", raw_output),
                    "prompt_final_path": logger.write_text(f"{slide_dir}/prompt_final.txt", final_prompt),
                    "prompt_fallback_used": not bool(raw_output),
                }
                logger.set_slide(slide.page, payload)
            return final_prompt
        except Exception as exc:
            if logger:
                fallback_prompt = self._ensure_prompt_density_guidance(
                    self._normalize_long_text_prompt(fallback),
                    information_density,
                )
                logger.set_slide(
                    slide.page,
                    {
                        "prompt_final_path": logger.write_text(
                            f"{slide_dir}/prompt_final.txt",
                            fallback_prompt,
                        ),
                        "prompt_fallback_used": True,
                        "prompt_error": str(exc),
                    },
                )
            return self._ensure_prompt_density_guidance(
                self._normalize_long_text_prompt(fallback),
                information_density,
            )

    def _build_page_prompt(
        self,
        deck_title: str,
        requirement: str,
        slide: SlideOutline,
        style_prompt: str,
        information_density: str = DEFAULT_INFORMATION_DENSITY,
    ) -> str:
        points = "；".join(x.strip() for x in slide.key_points if x.strip()) or "TBD point"
        normalized_density = self._normalize_information_density(information_density)
        density_guidance = self._build_information_density_guidance(normalized_density)
        density_section = (
            f"\n本页信息密度参考：\n{density_guidance}\n"
            if density_guidance
            else ""
        )
        return f"""
你是一名世界级演示视觉总监和图像生成提示词专家。现在请直接构造一段可用于生成整页 PPT 单页画面的完整提示词，这段提示词必须尽可能长、尽可能完整、尽可能具体、尽可能细致，宁可更详细也不要因为概括而漏掉细节。整页内容对应的全局上下文是：整套 PPT 标题为“{deck_title}”，用户需求是“{requirement}”，当前是第 {slide.page} 页，页面标题是“{slide.title}”，本页需要覆盖的核心要点包括：{points}。整页画面必须首先严格继承以下风格 DNA，并把它转化为当前页可执行、可落地、细到足以直接出图的画面描述：{style_prompt}
{density_section}

在这段完整提示词里，你要像写一份逐层施工说明一样，从整页的顶部到底部、从左到右，把所有重要的可见元素和布局关系都写清楚，不要省略关键视觉决策。要明确交代整页是完整宽屏 PPT 页面，能够完整承载页面结构与信息层级；所有页面内实际出现的文字都必须是简体中文；要完整保留本页标题和核心要点，不能遗漏内容模块。请把版头、标题区、副标题或说明文字、主体信息区、卡片或模块、图表或信息图、标签、连接线、箭头、边框、图标、辅助说明区、页码区、背景处理、强调色位置、阴影与材质感这些元素尽可能详细地描述出来，并且写清楚它们之间的相对位置、大小关系、对齐关系、主次层级、节奏变化和视觉流向，不要只给抽象概括，也不要简单说“放一个模块”。

这段提示词还必须充分说明排版气质和图解逻辑，包括页面骨架、留白方式、模块密度、信息分层、阅读动线、视觉锚点、对比关系、图解组织方式、卡片或容器处理、边框粗细倾向、分隔条样式、箭头语言、标签语法、图标风格、图表外观、光影和材质控制方式。要特别强调页面需要保持专业、清晰、适合正式汇报，整体可读性强，标题醒目，正文与标签自然，数字或重点结论有明确强调。如果本页天然适合流程图、对比图、分层图、时间线、矩阵图、结构图或指标模块，就要把相应的信息图表达方式写具体，而不是一句带过。整页视觉必须和整套 PPT 保持同一视觉家族，背景系统、色彩关系、卡片样式、边框语言、阴影质感、图标体系、图表语言、材质感和渲染气质都不能漂移；如果风格 DNA 本身暗示的是高密度技术图解、工程结构图、科研汇报式页面，或者更强装饰性、插画性、漫画性、拼贴性、极简性等鲜明表达，都应在一致性的前提下忠实延续，而不是被预设审美偏好改写。

请同时把与当前风格系统一致的约束自然地融入这段提示词中：不要让画面偏离、误读或削弱参考风格或既定风格 DNA，不要无理由改变视觉家族、材质逻辑、色温关系、组件语言、信息密度和渲染方式；只有当参考风格或风格 DNA 本身包含某种表达方式时，才保留该表达。不要在提示词里出现分辨率、像素、seed、steps、CFG、宽高比、输出格式、模型参数等技术字段。最终形成的应该是一段高度完整、细节充分、几乎不把关键视觉决策留给模型猜测的完整提示词。
""".strip()

    def _generate_slide_image(
        self,
        runtime_cfg: RuntimeConfig,
        prompt: str,
        output_path: Path,
        slide_page: Optional[int] = None,
        logger: Optional[GenerationLogger] = None,
    ) -> None:
        image_provider = (runtime_cfg.image_provider or "http").strip().lower()
        augmented_prompt = self._augment_prompt(prompt)
        slide_dir = f"slides/slide_{slide_page:02d}" if slide_page else ""
        if logger and slide_page is not None:
            image_prompt_path = logger.write_text(f"{slide_dir}/image_prompt.txt", augmented_prompt)
            image_request_path = logger.write_json(
                f"{slide_dir}/image_request.json",
                {
                    "provider": image_provider,
                    "image_base_url": runtime_cfg.image_base_url,
                    "model": runtime_cfg.image_model,
                    "image_size": self.settings.image_size,
                    "variants": self.settings.image_variants,
                    "timeout_seconds": self.settings.image_timeout,
                    "retries": self.settings.image_retries,
                    "image_prompt_path": image_prompt_path,
                },
            )
            logger.set_slide(
                slide_page,
                {
                    "image_prompt_path": image_prompt_path,
                    "image_request_path": image_request_path,
                },
            )

        if image_provider == "http":
            self._generate_http_slide_image(
                image_api_url=runtime_cfg.image_base_url,
                image_api_key=runtime_cfg.image_api_key,
                image_model=runtime_cfg.image_model,
                prompt=augmented_prompt,
                output_path=output_path,
                slide_page=slide_page,
                logger=logger,
            )
            return

        attempts = 1 + max(0, self.settings.image_retries)
        last_error: Optional[Exception] = None
        for attempt in range(1, attempts + 1):
            try:
                if image_provider == "openai":
                    image_bytes, response_summary = self._generate_openai_image_once(
                        image_base_url=runtime_cfg.image_base_url,
                        image_api_key=runtime_cfg.image_api_key,
                        image_model=runtime_cfg.image_model,
                        prompt=augmented_prompt,
                    )
                elif image_provider == "gemini":
                    image_bytes, response_summary = self._generate_gemini_image_once(
                        image_base_url=runtime_cfg.image_base_url,
                        image_api_key=runtime_cfg.image_api_key,
                        image_model=runtime_cfg.image_model,
                        prompt=augmented_prompt,
                    )
                else:
                    raise ValueError(f"Unsupported image provider: {image_provider}")

                output_path.parent.mkdir(parents=True, exist_ok=True)
                output_path.write_bytes(image_bytes)
                self._safe_enforce_widescreen_16x9(output_path)
                if logger and slide_page is not None:
                    logger.append_slide_event(
                        slide_page,
                        "image_attempts",
                        {
                            "attempt": attempt,
                            "success": True,
                            "provider": image_provider,
                            "response_summary": response_summary,
                        },
                    )
                return
            except Exception as exc:
                last_error = exc
                if logger and slide_page is not None:
                    logger.append_slide_event(
                        slide_page,
                        "image_attempts",
                        {
                            "attempt": attempt,
                            "success": False,
                            "provider": image_provider,
                            "error": str(exc),
                        },
                    )
                if attempt >= attempts:
                    break

        raise ValueError(f"image generation failed after {attempts} attempts: {last_error}")

    def _generate_http_slide_image(
        self,
        *,
        image_api_url: str,
        image_api_key: str,
        image_model: str,
        prompt: str,
        output_path: Path,
        slide_page: Optional[int],
        logger: Optional[GenerationLogger],
    ) -> None:
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {image_api_key}",
            "Accept": "application/json",
            "User-Agent": "ppt-image-generator/0.1",
        }
        payload = {
            "model": image_model,
            "prompt": prompt,
            "imageSize": self.settings.image_size,
            "variants": self.settings.image_variants,
            "shutProgress": True,
        }

        attempts = 1 + max(0, self.settings.image_retries)
        last_error: Optional[Exception] = None
        with requests.Session() as session:
            for attempt in range(1, attempts + 1):
                try:
                    resp = session.post(
                        image_api_url,
                        headers=headers,
                        json=payload,
                        timeout=self.settings.image_timeout,
                    )
                    raw = resp.text or ""
                    if not (200 <= resp.status_code < 300):
                        if logger and slide_page is not None:
                            logger.append_slide_event(
                                slide_page,
                                "image_attempts",
                                {
                                    "attempt": attempt,
                                    "success": False,
                                    "provider": "http",
                                    "status_code": resp.status_code,
                                    "response_excerpt": raw[:1000],
                                    "error": f"HTTP {resp.status_code}",
                                },
                            )
                        raise ValueError(f"HTTP {resp.status_code}: {raw[:300]}")

                    try:
                        data = resp.json()
                    except Exception:
                        data = self._relaxed_json_parse(raw)

                    moderation = self._moderation_reason(data, raw)
                    if moderation:
                        if logger and slide_page is not None:
                            logger.append_slide_event(
                                slide_page,
                                "image_attempts",
                                {
                                    "attempt": attempt,
                                    "success": False,
                                    "provider": "http",
                                    "status_code": resp.status_code,
                                    "response_summary": self._summarize_image_response(data, raw),
                                    "error": f"blocked by moderation: {moderation}",
                                },
                            )
                        raise ValueError(f"blocked by moderation: {moderation}")

                    urls = self._extract_urls(data)
                    if urls:
                        if logger and slide_page is not None:
                            logger.append_slide_event(
                                slide_page,
                                "image_attempts",
                                {
                                    "attempt": attempt,
                                    "success": True,
                                    "provider": "http",
                                    "status_code": resp.status_code,
                                    "download_url": urls[0],
                                    "response_summary": self._summarize_image_response(data, raw),
                                },
                            )
                        self._download_to_path(session, urls[0], output_path, self.settings.image_timeout)
                        self._safe_enforce_widescreen_16x9(output_path)
                        return

                    b64_json = str(data.get("b64_json") or "")
                    if b64_json:
                        if logger and slide_page is not None:
                            logger.append_slide_event(
                                slide_page,
                                "image_attempts",
                                {
                                    "attempt": attempt,
                                    "success": True,
                                    "provider": "http",
                                    "status_code": resp.status_code,
                                    "b64_length": len(b64_json),
                                    "response_summary": self._summarize_image_response(data, raw),
                                },
                            )
                        output_path.parent.mkdir(parents=True, exist_ok=True)
                        output_path.write_bytes(base64.b64decode(b64_json))
                        self._safe_enforce_widescreen_16x9(output_path)
                        return

                    if logger and slide_page is not None:
                        logger.append_slide_event(
                            slide_page,
                            "image_attempts",
                            {
                                "attempt": attempt,
                                "success": False,
                                "provider": "http",
                                "status_code": resp.status_code,
                                "response_excerpt": raw[:1000],
                                "error": "no urls in response",
                            },
                        )
                    raise ValueError(f"no urls in response: {raw[:300]}")
                except Exception as exc:
                    last_error = exc
                    if logger and slide_page is not None:
                        logger.append_slide_event(
                            slide_page,
                            "image_attempts",
                            {
                                "attempt": attempt,
                                "success": False,
                                "provider": "http",
                                "error": str(exc),
                            },
                        )
                    if attempt >= attempts:
                        break

        raise ValueError(f"image generation failed after {attempts} attempts: {last_error}")

    def _generate_openai_image_once(
        self,
        *,
        image_base_url: str,
        image_api_key: str,
        image_model: str,
        prompt: str,
    ) -> tuple[bytes, dict[str, Any]]:
        response = requests.post(
            self._build_openai_images_url(image_base_url),
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {image_api_key}",
            },
            json=self._build_openai_image_payload(image_model, prompt),
            timeout=self.settings.image_timeout,
        )
        raw = response.text or ""
        if not response.ok:
            raise ValueError(f"HTTP {response.status_code}: {raw[:300]}")
        try:
            data = response.json()
        except Exception as exc:
            raise ValueError(f"OpenAI image response was not JSON: {raw[:300]}") from exc

        image_bytes = self._extract_openai_image_bytes(data)
        summary = self._summarize_image_response(data, raw)
        summary["provider"] = "openai"
        return image_bytes, summary

    def _generate_gemini_image_once(
        self,
        *,
        image_base_url: str,
        image_api_key: str,
        image_model: str,
        prompt: str,
    ) -> tuple[bytes, dict[str, Any]]:
        response = requests.post(
            self._build_gemini_generate_content_url(image_base_url, image_model),
            params={"key": image_api_key},
            headers={"Content-Type": "application/json"},
            json={
                "contents": [
                    {
                        "role": "user",
                        "parts": [{"text": prompt}],
                    }
                ]
            },
            timeout=self.settings.image_timeout,
        )
        raw = response.text or ""
        if not response.ok:
            raise ValueError(f"HTTP {response.status_code}: {raw[:300]}")
        try:
            data = response.json()
        except Exception as exc:
            raise ValueError(f"Gemini image response was not JSON: {raw[:300]}") from exc

        self._raise_for_gemini_image_block(data)
        image_bytes, mime_type = self._extract_gemini_image_bytes(data)
        summary = {
            "provider": "gemini",
            "mime_type": mime_type,
            "keys": sorted(str(key) for key in data.keys()),
            "raw_excerpt": raw[:1000],
            "byte_length": len(image_bytes),
        }
        return image_bytes, summary

    def _build_pptx(self, slide_results: list[SlideResult], run_dir: Path, output_path: Path) -> None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        for item in sorted(slide_results, key=lambda x: x.page):
            image_path = run_dir / f"slide_{item.page:02d}.png"
            if not image_path.exists():
                raise ValueError(f"missing slide image for pptx: {image_path}")
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    @staticmethod
    def _summarize_image_response(data: Optional[dict[str, Any]], raw_text: str) -> dict[str, Any]:
        summary: dict[str, Any] = {
            "raw_excerpt": (raw_text or "")[:1000],
        }
        if not data:
            return summary
        summary["keys"] = sorted(str(key) for key in data.keys())
        results = data.get("results")
        if isinstance(results, list):
            summary["result_count"] = len(results)
        if data.get("url"):
            summary["has_top_level_url"] = True
        b64_json = data.get("b64_json")
        if b64_json:
            summary["b64_length"] = len(str(b64_json))
        return summary

    @staticmethod
    def _build_openai_images_url(base_url: str) -> str:
        cleaned = (base_url or "").strip().rstrip("/")
        if cleaned.endswith("/images/generations"):
            return cleaned
        return f"{cleaned}/images/generations"

    def _build_openai_image_payload(self, image_model: str, prompt: str) -> dict[str, Any]:
        payload: dict[str, Any] = {
            "model": image_model,
            "prompt": prompt,
            "n": self.settings.image_variants,
            "response_format": "b64_json",
        }
        if self.settings.image_size:
            payload["size"] = self.settings.image_size
        return payload

    def _extract_openai_image_bytes(self, payload: dict[str, Any]) -> bytes:
        data = payload.get("data")
        if isinstance(data, list):
            for item in data:
                if not isinstance(item, dict):
                    continue
                b64_json = item.get("b64_json")
                if b64_json:
                    try:
                        return base64.b64decode(str(b64_json))
                    except (binascii.Error, ValueError) as exc:
                        raise ValueError("OpenAI image response b64_json is invalid.") from exc
                image_url = item.get("url")
                if image_url:
                    response = requests.get(str(image_url), timeout=self.settings.image_timeout)
                    response.raise_for_status()
                    return response.content
        raise ValueError(f"OpenAI image response did not contain usable image data: {json.dumps(payload)[:500]}")

    @staticmethod
    def _build_gemini_generate_content_url(base_url: str, model: str) -> str:
        cleaned = (base_url or "").strip().rstrip("/")
        if cleaned.endswith(":generateContent"):
            return cleaned
        if cleaned.endswith("/models"):
            return f"{cleaned}/{model}:generateContent"
        if "/models/" in cleaned:
            return f"{cleaned}:generateContent"
        return f"{cleaned}/models/{model}:generateContent"

    @staticmethod
    def _raise_for_gemini_image_block(payload: dict[str, Any]) -> None:
        prompt_feedback = payload.get("promptFeedback") or {}
        block_reason = prompt_feedback.get("blockReason")
        if block_reason:
            raise ValueError(f"Gemini blocked the image request: {block_reason}")

        for candidate in payload.get("candidates") or []:
            if not isinstance(candidate, dict):
                continue
            finish_reason = candidate.get("finishReason")
            if finish_reason in {"SAFETY", "RECITATION", "BLOCKLIST", "PROHIBITED_CONTENT"}:
                raise ValueError(f"Gemini image request blocked with finish reason: {finish_reason}")

    @staticmethod
    def _extract_gemini_image_bytes(payload: dict[str, Any]) -> tuple[bytes, str]:
        for candidate in payload.get("candidates") or []:
            if not isinstance(candidate, dict):
                continue
            content = candidate.get("content") or {}
            for part in content.get("parts") or []:
                if not isinstance(part, dict):
                    continue
                inline_data = part.get("inlineData") or part.get("inline_data")
                if not isinstance(inline_data, dict):
                    continue
                encoded = inline_data.get("data")
                if not encoded:
                    continue
                mime_type = str(inline_data.get("mimeType") or inline_data.get("mime_type") or "image/png")
                try:
                    return base64.b64decode(str(encoded)), mime_type
                except (binascii.Error, ValueError) as exc:
                    raise ValueError("Gemini image response contained invalid base64 image data.") from exc
        raise ValueError(f"Gemini image response did not contain inline image data: {json.dumps(payload)[:500]}")

    @staticmethod
    def _mime_to_extension(mime: Optional[str]) -> str:
        mapping = {
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/jpg": ".jpg",
            "image/webp": ".webp",
        }
        return mapping.get((mime or "").strip().lower(), ".bin")

    @staticmethod
    def _image_bytes_to_data_url(image_bytes: Optional[bytes], mime: Optional[str]) -> Optional[str]:
        if not image_bytes:
            return None
        normalized_mime = (mime or "image/png").strip() or "image/png"
        return f"data:{normalized_mime};base64,{base64.b64encode(image_bytes).decode('utf-8')}"

    @staticmethod
    def _extract_urls(resp_json: dict[str, Any]) -> list[str]:
        results = resp_json.get("results") or []
        urls: list[str] = []
        if isinstance(results, list):
            for item in results:
                if isinstance(item, dict) and item.get("url"):
                    urls.append(str(item["url"]))
        if not urls and resp_json.get("url"):
            urls = [str(resp_json["url"])]
        return urls

    @staticmethod
    def _relaxed_json_parse(text: str) -> dict[str, Any]:
        raw = (text or "").strip()
        try:
            return json.loads(raw)
        except Exception:
            pass

        lines = [ln.strip() for ln in raw.splitlines() if ln.strip()]
        data_lines: list[str] = []
        for line in lines:
            if line.startswith("data:"):
                line = line[len("data:") :].strip()
            data_lines.append(line)
        merged = "\n".join(data_lines).strip()

        try:
            return json.loads(merged)
        except Exception:
            pass

        match = re.search(r"\{.*\}", merged, flags=re.DOTALL)
        if match:
            return json.loads(match.group(0))
        raise ValueError("unable to parse JSON from image API response")

    @staticmethod
    def _moderation_reason(data: Optional[dict[str, Any]], raw_text: str) -> Optional[str]:
        low = (raw_text or "").lower()
        if data:
            fr = str(data.get("failure_reason") or "").lower()
            if fr in ("input_moderation", "output_moderation"):
                return fr
            err = str(data.get("error") or "").lower()
            if "input_moderation" in err:
                return "input_moderation"
            if "output_moderation" in err:
                return "output_moderation"
        if "input_moderation" in low:
            return "input_moderation"
        if "output_moderation" in low:
            return "output_moderation"
        return None

    @staticmethod
    def _download_to_path(session: requests.Session, url: str, save_path: Path, timeout: int) -> None:
        with session.get(url, stream=True, timeout=timeout) as resp:
            resp.raise_for_status()
            save_path.parent.mkdir(parents=True, exist_ok=True)
            tmp = save_path.with_suffix(save_path.suffix + ".part")
            with open(tmp, "wb") as fp:
                for chunk in resp.iter_content(chunk_size=1024 * 128):
                    if chunk:
                        fp.write(chunk)
            os.replace(tmp, save_path)

    @staticmethod
    def _safe_enforce_widescreen_16x9(image_path: Path) -> None:
        try:
            PPTImagePipeline._enforce_widescreen_16x9(image_path)
        except Exception:
            return

    @staticmethod
    def _enforce_widescreen_16x9(image_path: Path) -> None:
        from PIL import Image  # type: ignore

        target_ratio = 16 / 9
        with Image.open(image_path) as opened:
            img = opened.copy()

        width, height = img.size
        if width <= 0 or height <= 0:
            return

        ratio = width / height
        if abs(ratio - target_ratio) / target_ratio <= 0.01:
            return

        if ratio > target_ratio:
            canvas_w = width
            canvas_h = int(round(width / target_ratio))
        else:
            canvas_h = height
            canvas_w = int(round(height * target_ratio))

        if canvas_w <= 0 or canvas_h <= 0:
            return

        if img.mode in ("RGBA", "LA"):
            background = Image.new("RGBA", (canvas_w, canvas_h), (255, 255, 255, 255))
            x = (canvas_w - width) // 2
            y = (canvas_h - height) // 2
            background.paste(img, (x, y), mask=img.split()[-1])
            background.convert("RGB").save(image_path, format="PNG")
            return

        background = Image.new("RGB", (canvas_w, canvas_h), (255, 255, 255))
        x = (canvas_w - width) // 2
        y = (canvas_h - height) // 2
        background.paste(img.convert("RGB"), (x, y))
        background.save(image_path, format="PNG")

    @staticmethod
    def _augment_prompt(prompt: str) -> str:
        cleaned = (prompt or "").strip()
        if not cleaned:
            return cleaned
        if PROMPT_SUFFIX.lower() in cleaned.lower():
            return cleaned
        return cleaned + "\n\n" + PROMPT_SUFFIX

    @staticmethod
    def _looks_like_prompt_heading(line: str) -> bool:
        stripped = (line or "").strip()
        if not stripped:
            return False
        if stripped.startswith("[") and stripped.endswith("]") and len(stripped) <= 40:
            return True
        if len(stripped) <= 24 and stripped.endswith(("：", ":")):
            return True
        return False

    @staticmethod
    def _remove_prompt_parameter_tokens(text: str) -> str:
        cleaned = STYLE_PROMPT_PARAMETER_RE.sub("", text or "")
        cleaned = re.sub(r"\b(?:PPT\s*单页|PPT页面)\b", "PPT单页", cleaned, flags=re.I)
        cleaned = re.sub(r"\s{2,}", " ", cleaned)
        cleaned = re.sub(r"([（(])\s*[，,、;；:：]?\s*([）)])", r"\1\2", cleaned)
        cleaned = re.sub(r"[，,、]\s*[，,、]", "，", cleaned)
        cleaned = re.sub(r"[:：]\s*$", "", cleaned)
        return cleaned.strip(" ,，;；")

    @staticmethod
    def _normalize_long_text_prompt(raw_text: str) -> str:
        text = (raw_text or "").strip()
        if not text:
            return ""
        if text.startswith("```"):
            parts = [part for part in text.split("```") if part.strip()]
            if parts:
                text = parts[0].strip()
                if "\n" in text and text.splitlines()[0].strip().lower() in {"text", "prompt", "markdown", "md"}:
                    text = "\n".join(text.splitlines()[1:]).strip()

        normalized_lines: list[str] = []
        for raw_line in text.splitlines():
            line = raw_line.strip()
            if not line:
                if normalized_lines and normalized_lines[-1] != "":
                    normalized_lines.append("")
                continue
            if ASSISTANT_META_LINE_RE.search(line):
                break
            if re.fullmatch(r"(?i)(prompt|final prompt|output|最终提示词|提示词|提示词正文)[:：]?", line):
                continue
            if re.match(r"^#{1,6}\s*", line):
                heading = re.sub(r"^#{1,6}\s*", "", line).strip()
                if heading:
                    normalized_lines.append(f"[{heading}]")
                continue
            line = re.sub(r"\s+", " ", line).strip()
            if not line:
                continue
            if re.fullmatch(r"(?:[-*•·●]|\d+[\.、])\s*", line):
                continue
            if line.startswith(("-", "*", "•", "·", "●")):
                line = re.sub(r"^[-*•·●]\s*", "- ", line)
            elif re.match(r"^\d+[\.、]\s*", line):
                line = re.sub(r"^\d+[\.、]\s*", "- ", line)
            line = PPTImagePipeline._remove_prompt_parameter_tokens(line)
            if not line:
                continue
            if not PPTImagePipeline._looks_like_prompt_heading(line) and not re.search(r"[。！？；;：:，,\.\!]$", line):
                line += "。"
            normalized_lines.append(line)

        compact_lines: list[str] = []
        for line in normalized_lines:
            if line == "" and (not compact_lines or compact_lines[-1] == ""):
                continue
            compact_lines.append(line)
        return "\n".join(compact_lines).strip()

    @staticmethod
    def _parse_json_object(raw_text: str) -> dict[str, Any]:
        text = raw_text.strip()
        if text.startswith("```"):
            parts = [p for p in text.split("```") if p.strip()]
            if parts:
                text = parts[0]
                if text.lower().startswith("json"):
                    text = text[4:].strip()
        left = text.find("{")
        right = text.rfind("}")
        if left < 0 or right < 0 or right <= left:
            raise ValueError("Model did not return valid JSON.")
        return json.loads(text[left : right + 1])

    @staticmethod
    def _message_text(content: Any) -> str:
        if content is None:
            return ""
        if isinstance(content, str):
            return content
        if isinstance(content, Iterable):
            chunks: list[str] = []
            for item in content:
                if isinstance(item, str):
                    chunks.append(item)
                    continue
                if isinstance(item, dict):
                    text_val = item.get("text")
                    if text_val:
                        chunks.append(str(text_val))
                    continue
                text_val = getattr(item, "text", None)
                if text_val:
                    chunks.append(str(text_val))
            return "\n".join(chunks)
        return str(content)
